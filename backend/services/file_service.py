import json
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image

from services.llm_router import ask_llm


def extract_pdf_text(file_path: str):
    reader = PdfReader(file_path)
    text = []

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text.append(page_text)

    return "\n".join(text).strip()


def extract_docx_text(file_path: str):
    doc = Document(file_path)
    text = []

    for para in doc.paragraphs:
        if para.text.strip():
            text.append(para.text)

    return "\n".join(text).strip()


def extract_txt_text(file_path: str):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()


def extract_csv_text(file_path: str):
    df = pd.read_csv(file_path)
    return df.to_string(index=False)


def extract_excel_text(file_path: str):
    excel = pd.read_excel(file_path, sheet_name=None)
    text = []

    for sheet_name, df in excel.items():
        text.append(f"\nSheet: {sheet_name}\n")
        text.append(df.to_string(index=False))

    return "\n".join(text).strip()


def extract_pptx_text(file_path: str):
    prs = Presentation(file_path)
    text = []

    for i, slide in enumerate(prs.slides, start=1):
        text.append(f"\nSlide {i}\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)

    return "\n".join(text).strip()


def extract_json_text(file_path: str):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        data = json.load(f)

    return json.dumps(data, indent=2)


async def extract_image_text(file_path: str):
    prompt = """
Extract all visible text from this image.
Also explain the image content clearly.
"""

    # Gemini vision through llm_router
    return await ask_llm(prompt)


async def extract_file_content(file_path: str, extension: str):
    extension = extension.lower()

    if extension == "pdf":
        return extract_pdf_text(file_path), "pdf"

    if extension == "docx":
        return extract_docx_text(file_path), "docx"

    if extension == "txt":
        return extract_txt_text(file_path), "txt"

    if extension == "csv":
        return extract_csv_text(file_path), "csv"

    if extension in ["xlsx", "xls"]:
        return extract_excel_text(file_path), "excel"

    if extension == "pptx":
        return extract_pptx_text(file_path), "pptx"

    if extension == "json":
        return extract_json_text(file_path), "json"

    if extension in ["png", "jpg", "jpeg"]:
        return await extract_image_text(file_path), "image"

    return "", "unsupported"